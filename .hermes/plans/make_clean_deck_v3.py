#!/usr/bin/env python3
"""Clean, white Milestone 1 findings deck.
Images are embedded directly from primary-source PDF/page renders in source_shots/.
No highlight overlays or synthetic marks are added to screenshots.
"""
import os, json, re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = Path(r"C:\Dev\CyberZen\.hermes\plans")
SHOT_DIR = ROOT / "source_shots_v3"
MANIFEST = json.loads((SHOT_DIR / "manifest.json").read_text(encoding="utf-8"))
SOURCE = {(x["key"][0], x["key"][1]): x for x in MANIFEST}

# Reuse the researched finding titles, references, and comments already prepared in the first deck.
old_script = (ROOT / "make_deck2.py").read_text(encoding="utf-8")
prefix = old_script.split("# =====================================================================\n# BUILD")[0]
ns = {}
exec(prefix, ns)
SLIDES = ns["slides"]

# Source-grounded wording changes where the original Word synthesis was broader than the primary excerpt.
TITLE_OVERRIDES = {
    ("FA", "Framework gap · The central finding"): "The official pipeline starts at TRL 1 — but scale-up begins later",
    ("FA", "CRL · No published methodology"): "CRL: the commercial dimension TRL misses",
    ("FA", "Innosuisse · Custom, not codified"): "Innosuisse: bespoke, voucher-based coaching",
    ("FA", "Innosuisse · The TRL 1–3 structural gap"): "Innosuisse support is stage-specific",
    ("FA", "EIC · The accelerator gap"): "EIC Accelerator: TRL 6–8",
    ("DC", "Data · TTO survival: p=0.0034"): "TTOs become playmakers through relationships",
}
# Insight-layer comments v2: every bullet must add something the title/excerpt
# does not say — interpretation, cross-slide link, or open question.
import importlib.util as _ilu
_spec = _ilu.spec_from_file_location("comments_v2", ROOT / "comments_v2.py")
_mod = _ilu.module_from_spec(_spec); _spec.loader.exec_module(_mod)
COMMENT_OVERRIDES = _mod.NEW_COMMENTS

# White, restrained palette.
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x16, 0x1B, 0x22)
BODY = RGBColor(0x3F, 0x4A, 0x59)
MUTED = RGBColor(0x6B, 0x76, 0x86)
RULE = RGBColor(0xD9, 0xDF, 0xE7)
BLUE = RGBColor(0x1D, 0x4E, 0x89)
TEAL = RGBColor(0x0F, 0x8A, 0x82)
PURPLE = RGBColor(0x6D, 0x4A, 0xA5)
SKY = RGBColor(0x1C, 0x76, 0xA8)

ACCENT = {"LIT": TEAL, "FA": PURPLE, "DC": SKY}
SECTION = {"LIT": "LITERATURE REVIEW", "FA": "FRAMEWORK ASSESSMENT", "DC": "DATA COMPENDIUM"}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 42


def bg_slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = WHITE
    r.line.fill.background(); r.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for content, size, bold, color, after in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align; p.space_after = Pt(after)
        rr = p.add_run(); rr.text = content
        rr.font.name = "Arial"; rr.font.size = Pt(size); rr.font.bold = bold; rr.font.color.rgb = color
    return tb


def line(s, x, y, w, color=RULE, height=0.012):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(height))
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background(); r.shadow.inherit = False


def section_header(s, tag, accent, page_no):
    text(s, 0.62, 0.35, 5.4, 0.25, [("NCCR MUONIVERSE  ·  HSG", 9, True, MUTED, 0)])
    text(s, 0.62, 0.68, 3.8, 0.25, [(tag, 9.5, True, accent, 0)])
    text(s, 12.0, 0.35, 0.7, 0.25, [(f"{page_no:02d}", 9.5, True, MUTED, 0)], align=PP_ALIGN.RIGHT)
    line(s, 0.62, 1.02, 12.08)


def divider(tag, accent, title, subtitle, page_no):
    s = bg_slide()
    section_header(s, tag, accent, page_no)
    # simple centered section page
    text(s, 0.82, 3.05, 11.7, 0.55, [(title, 34, True, INK, 0)])
    text(s, 0.84, 3.78, 10.5, 0.35, [(subtitle, 14, False, MUTED, 0)])
    line(s, 0.84, 4.48, 1.2, accent, 0.035)
    text(s, 0.84, 6.82, 10, 0.25, [("Milestone 1 · source-led findings deck", 9.5, False, MUTED, 0)])
    return s


def source_key(key):
    return SOURCE[key]


def finding(tag, key, title, old_ref, old_comments, page_no):
    s = bg_slide(); accent = ACCENT[tag]
    section_header(s, SECTION[tag], accent, page_no)
    title = TITLE_OVERRIDES.get(key, title)
    comments = COMMENT_OVERRIDES.get(key, old_comments)
    src = source_key(key)

    # Title and short descriptor.
    text(s, 0.62, 1.30, 8.0, 0.30, [("SOURCE-LED FINDING", 9.5, True, accent, 0)])
    text(s, 0.62, 1.62, 12.0, 0.70, [(title, 25, True, INK, 0)])

    # Original-source excerpt, left side. This image is untouched after rendering from the cited PDF.
    img_path = str(SHOT_DIR / src["file_out"])
    with Image.open(img_path) as im:
        iw, ih = im.size
    max_w, max_h = 7.15, 3.75
    scale = min(max_w / iw, max_h / ih)
    dw, dh = iw * scale, ih * scale
    ix, iy = 0.62, 2.55
    s.shapes.add_picture(img_path, Inches(ix), Inches(iy), Inches(dw), Inches(dh))
    frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ix - 0.025), Inches(iy - 0.025), Inches(dw + 0.05), Inches(dh + 0.05))
    frame.fill.background(); frame.line.color.rgb = RULE; frame.line.width = Pt(0.8); frame.shadow.inherit = False

    text(s, 0.62, 2.28, 4.4, 0.20, [(f"ORIGINAL SOURCE EXCERPT  ·  p. {src['page']}", 8.5, True, MUTED, 0)])
    ref_y = iy + dh + 0.24
    text(s, 0.62, ref_y, 7.15, 0.62, [(src["ref"], 9.2, False, MUTED, 0)])

    # Right column: comments only, no overlay on image.
    rx, ry, rw = 8.18, 2.28, 4.48
    line(s, rx, ry, 0.55, accent, 0.035)
    text(s, rx, ry + 0.24, rw, 0.25, [("COMMENTS", 10, True, INK, 0)])
    cy = ry + 0.66
    for c in comments:
        est_lines = max(2, -(-len(c) // 56))
        h = est_lines * 0.185
        # a very small rule marker, not a bullet glyph
        line(s, rx, cy + 0.08, 0.15, accent, 0.018)
        text(s, rx + 0.28, cy, rw - 0.28, h, [(c, 10.5, False, BODY, 0)])
        cy += h + 0.16

    # Footer, simple and unobtrusive.
    line(s, 0.62, 7.03, 12.08)
    text(s, 0.62, 7.13, 7.0, 0.18, [("Lorik Dalloshi · Milestone 1 · NCCR Muoniverse", 8.5, False, MUTED, 0)])
    text(s, 11.4, 7.13, 1.3, 0.18, [(f"{page_no:02d} / {TOTAL:02d}", 8.5, True, MUTED, 0)], align=PP_ALIGN.RIGHT)
    return s

# Title slide.
s = bg_slide()
text(s, 0.72, 0.62, 5.5, 0.25, [("NCCR MUONIVERSE  ·  HSG", 10, True, TEAL, 0)])
line(s, 0.72, 1.22, 1.35, TEAL, 0.045)
text(s, 0.72, 2.08, 11.5, 0.9, [("Milestone 1 Findings", 42, True, INK, 3)])
text(s, 0.72, 3.10, 10.8, 0.45, [("Literature review · framework assessment · data compendium", 18, False, BODY, 0)])
text(s, 0.72, 4.58, 8.3, 0.75, [("A clean source-led deck: every finding uses an excerpt rendered directly from the cited paper, report, or official page.", 13, False, MUTED, 0)])
text(s, 0.72, 6.75, 8.0, 0.22, [("Prepared for the NCCR Muoniverse team · August 2026", 9.5, False, MUTED, 0)])

page = 2
# Literature
page += 0; divider(SECTION["LIT"], ACCENT["LIT"], "Seven themes", "What the primary literature converges on", page); page += 1
for kind,key,title,ref,comments in SLIDES:
    if kind == "theme":
        finding("LIT", key, title, ref, comments, page); page += 1
# Frameworks
divider(SECTION["FA"], ACCENT["FA"], "Four frameworks, one gap", "RCA / Imperial · CRL · Innosuisse · EIC · candidates", page); page += 1
for kind,key,title,ref,comments in SLIDES:
    if kind == "fw":
        finding("FA", key, title, ref, comments, page); page += 1
# Data
divider(SECTION["DC"], ACCENT["DC"], "The numbers", "Statistical evidence behind the findings", page); page += 1
for kind,key,title,ref,comments in SLIDES:
    if kind == "dt":
        finding("DC", key, title, ref, comments, page); page += 1

# Closing.
s = bg_slide(); section_header(s, "MILESTONE 1", TEAL, page)
text(s, 0.72, 2.65, 11.5, 0.65, [("Source-led evidence → interview questions", 31, True, INK, 0)])
text(s, 0.72, 3.55, 10.5, 0.9, [("The practical conclusion: early-TRL physics commercialization is not missing ideas — it is missing a support architecture designed for the gap.", 15, False, BODY, 4), ("Next: Milestone 2 interviews · September 1, 2026", 12, True, TEAL, 0)])
line(s, 0.72, 5.15, 1.3, TEAL, 0.04)
text(s, 0.72, 6.75, 10, 0.22, [("All source excerpts are embedded from the primary PDFs/pages listed beneath each image.", 9.5, False, MUTED, 0)])

out = Path(r"C:\Dev\CyberZen\.hermes\plans\Milestone1_Findings_Deck_accurate_v3.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
# Preserve the earlier dark version before replacing it.
backup = out.with_name("Milestone1_Findings_Deck_wallet_backup.pptx")
if out.exists() and not backup.exists():
    out.replace(backup)
prs.save(str(out))
print("SAVED", out, "slides", len(prs.slides), "bytes", out.stat().st_size)
print("SOURCE IMAGES", len(SOURCE))
