#!/usr/bin/env python3
"""Build the Milestone 1 Findings deck — wallet-style layout:
title top, highlighted doc screenshot left, reference under it, comments right."""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

SHOTS = r"C:\Dev\CyberZen\.hermes\plans\shots"
MANIFEST = json.load(open(os.path.join(SHOTS, "manifest.json"), encoding="utf-8"))
by_label = {(m["doc"], m["label"]): m for m in MANIFEST}

# ---------------- palette ----------------
BG      = RGBColor(0x0B, 0x12, 0x20)
BAR     = RGBColor(0x11, 0x1C, 0x2E)
CARD    = RGBColor(0x15, 0x22, 0x38)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SLATE   = RGBColor(0xC9, 0xD6, 0xE3)
GRAY    = RGBColor(0x6B, 0x7F, 0x96)
TEAL    = RGBColor(0x14, 0xB8, 0xA6)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
VIOLET  = RGBColor(0x8B, 0x5C, 0xF6)
SKY     = RGBColor(0x38, 0xBD, 0xF8)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
RED     = RGBColor(0xEF, 0x44, 0x44)

FONT = "Calibri"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s

def textbox(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for (text, size, bold, color, space) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space)
        r = p.add_run(); r.text = text
        f = r.font; f.name = FONT; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    return tb

def chip(s, x, y, w, h, text, fill=TEAL, txt=WHITE, size=10):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.06); tf.margin_top = tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = txt
    return c

def wallet_bar(s, section, accent, page_no, total):
    """Top bar mimicking the wallet screenshot: back arrow, title, three dots."""
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.62))
    bar.fill.solid(); bar.fill.fore_color.rgb = BAR
    bar.line.fill.background(); bar.shadow.inherit = False
    # back arrow (chevron)
    textbox(s, 0.35, 0.14, 0.4, 0.4, [("\u2039", 22, True, WHITE, 0)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, 0.75, 0.14, 6.5, 0.4, [("Milestone 1  \u00b7  Findings", 13, True, WHITE, 0)], anchor=MSO_ANCHOR.MIDDLE)
    # three dots right
    dots = "\u00b7  \u00b7  \u00b7"
    textbox(s, 12.15, 0.14, 0.9, 0.4, [(dots, 16, True, GRAY, 0)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # section chip under bar
    chip(s, 0.55, 0.78, 1.9, 0.34, section, fill=accent, size=9.5)

def finding_slide(section, accent, title, kicker, shot_key, reference, comments, page_no, total, big_ref=False):
    s = add_slide()
    wallet_bar(s, section, accent, page_no, total)
    # kicker + title
    textbox(s, 0.55, 1.28, 12.2, 0.35, [(kicker.upper(), 10.5, True, accent, 0)])
    textbox(s, 0.55, 1.58, 12.2, 0.62, [(title, 23, True, WHITE, 0)])
    # ------- left: screenshot -------
    shot = by_label[shot_key]
    img_path = os.path.join(SHOTS, shot["file"])
    im = Image.open(img_path)
    iw, ih = im.size
    max_w, max_h = 7.25, 3.95  # leave room for reference below
    ratio = min(max_w / iw, max_h / ih)
    disp_w, disp_h = iw * ratio, ih * ratio
    img_x, img_y = 0.55, 2.28
    pic = s.shapes.add_picture(img_path, Inches(img_x), Inches(img_y), Inches(disp_w), Inches(disp_h))
    # frame border
    frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(img_x - 0.06), Inches(img_y - 0.06),
                               Inches(disp_w + 0.12), Inches(disp_h + 0.12))
    frame.fill.background()
    frame.line.color.rgb = RGBColor(0x2A, 0x3A, 0x55); frame.line.width = Pt(1)
    frame.shadow.inherit = False
    # move frame behind picture
    sp = frame._element
    sp.getparent().remove(sp)
    pic._element.addprevious(sp)
    # quote caption chip on top of screenshot
    chip(s, img_x + 0.12, img_y + 0.12, 2.6, 0.32, "QUOTE HIGHLIGHTED", fill=accent, size=8.5)
    # ------- reference under screenshot -------
    ref_y = img_y + disp_h + 0.28
    chip(s, 0.55, ref_y, 1.05, 0.3, "SOURCE", fill=RGBColor(0x1E, 0x3A, 0x5F), size=8.5)
    textbox(s, 1.75, ref_y - 0.02, 6.0, 0.9, [(reference, 9.5 if not big_ref else 9, False, GRAY, 0)])
    # ------- right: comments card -------
    card_x, card_y = 8.25, 2.28
    card_w, card_h = 4.55, 4.55
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(card_x), Inches(card_y), Inches(card_w), Inches(card_h))
    card.fill.solid(); card.fill.fore_color.rgb = CARD
    card.line.color.rgb = RGBColor(0x2A, 0x3A, 0x55); card.line.width = Pt(1)
    card.shadow.inherit = False
    textbox(s, card_x + 0.3, card_y + 0.22, card_w - 0.6, 0.35, [("COMMENTS", 10.5, True, accent, 0)])
    tb_y = card_y + 0.62
    bullets = []
    for c in comments:
        bullets.append((c, 11, False, SLATE, 8))
    textbox(s, card_x + 0.3, tb_y, card_w - 0.62, card_h - 0.85, bullets)
    # bottom status row (wallet-like balance strip)
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.05), prs.slide_width, Inches(0.45))
    strip.fill.solid(); strip.fill.fore_color.rgb = BAR
    strip.line.fill.background(); strip.shadow.inherit = False
    textbox(s, 0.55, 7.12, 6.0, 0.3, [("NCCR Muoniverse \u00b7 HSG \u00b7 Milestone 1 findings", 9, False, GRAY, 0)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, 11.6, 7.12, 1.2, 0.3, [(f"{page_no:02d} / {total:02d}", 9.5, True, GRAY, 0)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def divider(section, accent, title, subtitle, page_no, total):
    s = add_slide()
    wallet_bar(s, section, accent, page_no, total)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(3.1), Inches(0.09), Inches(1.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background(); bar.shadow.inherit = False
    textbox(s, 0.9, 3.15, 11.5, 0.7, [(title, 34, True, WHITE, 0)])
    textbox(s, 0.9, 3.95, 11.5, 0.5, [(subtitle, 14, False, GRAY, 0)])
    textbox(s, 0.9, 6.3, 11.5, 0.4, [(f"{page_no:02d} / {total:02d}", 9.5, True, GRAY, 0)])

# =====================================================================
# CONTENT
# =====================================================================
total = 40  # 3 dividers + 37 findings

# ---- SECTION 1: THEMES ----
S = "LITERATURE REVIEW"
A = TEAL
slides = []
def th(key, title, ref, comments):
    slides.append(("theme", key, title, ref, comments))

th(("LIT", "Theme 0 · The headline finding"),
   "No existing framework was designed for TRL 1\u20133",
   "Milestone 1 Literature Review (2026), Executive Summary, p. 1.",
   ["The space where NCCR Muoniverse operates \u2014 fundamental physics at TRL 1\u20133 \u2014 is precisely where no commercialization framework was designed to work.",
    "Every framework assessed (RCA/Imperial, CRL, Innosuisse, EIC) has a structural gap at this stage.",
    "This gap is the thesis's core motivation: identify what support mechanisms fill it."])

th(("LIT", "Theme 1 · Deep tech is ecosystem-dependent"),
   "Ecosystem, not policy, drives deep tech",
   "Raff-Heinen et al. (2026), J. Technology Transfer \u00b7 Seitz et al. (2026), Review of Managerial Science.",
   ["Knowledge infrastructure (IRR 1.48) and talent pools (IRR 1.66) predict deep tech concentration; formal institutions are NOT significant.",
    "Top 10 EU regions hold 40% of all deep tech startups \u2014 concentration wins.",
    "University-affiliated incubators (r = 0.118) outperform generic ones (r = 0.067); mentoring is slightly negative (r = \u22120.012).",
    "\u2192 Switzerland's PSI/ETH/EPFL cluster is the real differentiator, not policy frameworks."])

th(("LIT", "Theme 2 · TTOs: relationships, not transactions"),
   "TTOs: relationships beat transactions",
   "Kruachottikul et al. (2023), J. Technology Transfer (231 spin-offs, 15-yr panel) \u00b7 Chen et al. (2024) \u00b7 Mart\u00ednez-Ardila et al. (2023).",
   ["Survival predictors, by significance: business relationships (p = 0.0034) > customer intros (p = 0.006) > market knowledge (p = 0.015) > tech knowledge (p = 0.042).",
    "'Playmaker' TTOs (active orchestrators) outperform 'go-between' TTOs.",
    "TTO specialization: depth coeff. 1.273 vs breadth 0.035 \u2014 depth wins.",
    "\u2192 PSI's TTO should adopt a playmaker orientation."])

th(("LIT", "Theme 3 · Structured frameworks outperform ad hoc"),
   "Structured frameworks work \u2014 if adapted",
   "Hayter et al. (2025), J. Technology Transfer \u00b7 VentureWell (2019) \u00b7 Barron & Amor\u00f3s (2020).",
   ["I-Corps (lean + customer discovery, 7-week immersion) independently adopted in 8 countries \u2014 a natural experiment: structured works, but every system adapted it locally.",
    "Frameworks need explicit abandon/pivot decision points \u2014 sunk-cost protection on 5\u201315 yr timelines.",
    "100 customer interviews in 9 weeks = paradigm shift in scientists' market understanding."])

th(("LIT", "Theme 4 · Support must be configured, not provided"),
   "Configuring support matters more than providing it",
   "Stahl et al. (2023) \u00b7 Hello Tomorrow (2019) \u00b7 Seitz et al. (2026).",
   ["3 effective configurations of 5 orchestration functions \u2014 no single configuration is universally optimal.",
    "Founders value mentors (64%) and customer intros (61%) above investment (37%) \u2014 capital is NOT Europe's scarcest resource.",
    "Networking helps (r = 0.071); more mentoring can even reduce performance (opportunity cost)."])

th(("LIT", "Theme 5 · Market identification is the binding constraint"),
   "Market is the binding constraint",
   "Foresight (2024), J. Innovation & Knowledge (n = 110 EU deep tech firms) \u00b7 Cantner et al. (2024).",
   ["Market features (β = 0.514) predict success nearly 2\u00d7 better than technology features (β = 0.310).",
    "Firm internal factors (team quality, management) are NOT significant.",
    "Cantner (n = 1,149): the commercial sphere dominates from day one of spin-off creation.",
    "\u2192 For muons: finding applications beats scientific excellence."])

th(("LIT", "Theme 5b · Cantner: commercial from day one"),
   "Commercial exposure from day one",
   "Cantner, Doerr, Goethner et al. (2024), Small Business Economics 62:1555\u20131590 (n = 1,149 scientists).",
   ["Dominance analysis: commercial sphere (industry contacts, market knowledge) dominates from the very outset \u2014 not in later stages.",
    "Challenges the assumption 'research first, commercialize later' \u2014 empirically wrong.",
    "\u2192 Support design principle: start commercial exposure immediately."])

th(("LIT", "Theme 5c · Denoo: business advisors help, tech advisors delay"),
   "Advisor type matters more than quantity",
   "Denoo, Van Boxstael & Belz (2024), J. Technology Transfer 49:1567\u20131705 (n = 112 ventures).",
   ["Each business advisor: +0.23 application-readiness units (p = 0.04).",
    "Each technology advisor: \u22120.31 units (p = 0.05) \u2014 actively delays.",
    "4 business advisors \u2248 +1 full unit \u2248 15 years of entrepreneurial experience.",
    "\u2192 Counterintuitive but evidence-based: limit tech advisors, prioritize business advisors."])

th(("LIT", "Theme 6 · Europe's structural barriers"),
   "Europe's structural barriers",
   "Stryber (2025) \u00b7 UNDP (2025).",
   ["€58B EU deep tech VC vs €215B US \u2014 a 3.7\u00d7 gap.",
    "62% of EU investors reject failed founders vs 23% US; 15% vs 45% of PhDs pursue startups.",
    "12% vs 40% of corporations engage with deep tech.",
    "\u2192 Cultural and structural \u2014 cannot be fixed by capital alone."])

th(("LIT", "Theme 6b · UNDP: 35% more time, 48% more capital"),
   "Deep tech needs 35% more time, 48% more capital",
   "UNDP (2025), Global Deep Tech Ecosystems.",
   ["Deep tech takes 35% more time and 48% more capital to reach equivalent milestones.",
    "These compound with European structural deficits: 3\u20135 yr VC fund cycles clash with deep tech timelines.",
    "\u2192 Muon tech's 10\u201320 yr horizon exceeds every Swiss instrument's design parameters."])

th(("LIT", "Theme 7 · Convergence on a deep tech definition"),
   "Converging on what deep tech is",
   "Kortsch et al. (2024), MIT Sloan / Univ. M\u00fcnster \u00b7 Cabanes et al. (2026), IEEE EMR.",
   ["12 defining attributes across 3 levels: technology (science-based, hardware, platforms), venture (5\u201315+ yr timelines, capital intensity), ecosystem (knowledge, talent, infrastructure).",
    "Cabanes: deep tech as a 'performative innovation regime' \u2014 patient capital, long-term partnerships, ecosystem orchestration.",
    "\u2192 Definitional convergence legitimizes a physics-specific commercialization lens."])

# ---- SECTION 2: FRAMEWORKS ----
S2 = "FRAMEWORK ASSESSMENT"
A2 = VIOLET
def fw(key, title, ref, comments):
    slides.append(("fw", key, title, ref, comments))

fw(("FA", "Framework gap · The central finding"),
   "Every framework gaps at TRL 1\u20133",
   "Milestone 1 Framework Assessment V3 (2026), Executive Summary.",
   ["RCA/Imperial: designed for TRL 4\u20136. CRL: conceptually 1\u20139 but operationally thin. Innosuisse: optimized for TRL 4\u20137. EIC Accelerator: TRL 5+.",
    "The space where Muoniverse operates is where no framework was designed to work.",
    "\u2192 Thesis contribution: adapt/combine existing frameworks to reach this unaddressed space."])

fw(("FA", "RCA/Imperial · What the module actually is"),
   "RCA/Imperial: a cluster, not a program",
   "Imperial College Business School; RCA; Imperial Innovation Challenge (student blog).",
   ["4 components: Innovation Challenge (1-week MBA), Wicked Module (1-month, with Steve Blank 2023), Wicked Acceleration Labs (WAL), DT Prime (£50\u2013500K non-dilutive venture builder).",
    "The Challenge is a pedagogical exercise, not a spin-off mechanism.",
    "\u2192 Methodology transferable; structure is not (institution-bound, MBA/design students)."])

fw(("FA", "RCA/Imperial · Flipped Lean Startup"),
   "RCA/Imperial: 'flipping' Lean Startup",
   "Pinder, M. (2021), mikepinder.co.uk \u00b7 Blank, S. (2023), 'Lean Meets Wicked Problems'.",
   ["Standard Lean Startup starts with a customer problem; deep tech starts with a technology capability and often NO existing market.",
    "WAL inverts the logic: technology-first speculative design \u2192 then validate with lead/extreme users.",
    "Endorsed by Steve Blank as a legitimate evolution for deep tech.",
    "\u2192 The key asset NCCR Muoniverse should adopt: the flipped methodology."])

fw(("FA", "CRL · The 9-point scale"),
   "CRL: the 9-point commercial scale",
   "Abbas & Nomvar, Commercial Readiness Level framework, U. Sydney / Scimita Ventures.",
   ["Runs parallel to NASA TRL; with MRL forms three 'orthogonal' maturity axes.",
    "Key concept: the binding constraint \u2014 whichever axis lags blocks progress. Rule: CRL should always \u2265 TRL.",
    "CRL 1\u20133 covers opportunity recognition \u2192 market sizing \u2192 competitiveness \u2014 conceptually relevant to TRL 1\u20133."])

fw(("FA", "CRL · The documentation problem"),
   "CRL: the documentation problem",
   "Abbas & Nomvar (n.d.), U. Sydney / Scimita Ventures \u2014 no journal, no DOI.",
   ["Despite being one of the most-cited 'commercial readiness' concepts, CRL has NO peer-reviewed primary paper.",
    "Exists as proprietary consultancy IP; cited by Deep Tech Leaders, ITONICS, NITI Aayog.",
    "The '22_CRL_Abbas.pdf' in the evidence folder is a DIFFERENT EIT Urban Mobility variant \u2014 not Abbas & Nomvar.",
    "\u2192 Must be cited as a practitioner framework, not an academically validated instrument."])

fw(("FA", "CRL · No published methodology"),
   "CRL: no published assessment methodology",
   "Abbas & Nomvar (n.d.) \u2014 Scimita Ventures proprietary service.",
   ["No per-level evidence checklists, no rubric, no independent assessment protocol, no certification.",
    "Descriptive (where you are) not prescriptive (what to do next) \u2014 unlike IRL/VIRAL.",
    "Operationally thin at low levels: CRL 1 looks identical for a muon venture and a mobile app.",
    "\u2192 Best used as a diagnostic framing tool (a detector can be TRL 9 but CRL 1)."])

fw(("FA", "Innosuisse · The native framework"),
   "Innosuisse: the native Swiss framework",
   "Innosuisse (SAFIG Act, SR 420.2); innosuisse.admin.ch.",
   ["~CHF 300M/yr, ~20 instruments \u2014 the most institutionally relevant framework: any NCCR venture will interact with it.",
    "Coaching: Initial (CHF 10K/12 mo) \u2192 Core (CHF 50K/36 mo, milestone-gated) \u2192 Scale-up (CHF 75K/24 mo).",
    "Voucher system: startups choose their own accredited coaches."])

fw(("FA", "Innosuisse · Custom, not codified"),
   "Innosuisse: custom, not codified",
   "Innosuisse Start-up Coaching; venturelab Entrepreneurship Training.",
   ["Coaching methodology is NOT explicitly based on I-Corps or a formal Lean Startup curriculum.",
    "Mentor-driven, bespoke per startup \u2014 fundamentally different from I-Corps' standardized 7-week immersion.",
    "Entrepreneurship Training (venturelab, since 2004) IS curriculum-based \u2014 alumni include Climeworks, GetYourGuide.",
    "\u2192 A direct methodological contrast for the benchmarking: mentor-driven vs methodology-driven."])

fw(("FA", "Innosuisse · KOF causal evidence"),
   "Innosuisse: KOF-verified causal effect",
   "Hulfeld, Spescha & W\u00f6rter (KOF ETH Zurich, Nov 2024), DiD study.",
   ["Only independently verified causal estimate: +20.7% sales, +17.6% employment over 5 years (matched control of 537 firms).",
    "Effects accelerate: employment +7.8% (yr 0) \u2192 +17.2% (yr 2) \u2192 +27.9% (yr 4).",
    "Strongest for small firms (5\u201350 employees) with ETH Domain research partners.",
    "\u2192 The only framework in the benchmark with causal validation."])

fw(("FA", "Innosuisse · The TRL 1–3 structural gap"),
   "Innosuisse: no TRL 1\u20133 instrument",
   "Framework Assessment V3, \u00a73.6 Structural Gap.",
   ["No instrument designed for TRL 1\u20133 commercialization; BRIDGE PoC (12 months, CHF 130K) is too short for muon validation cycles.",
    "Coaching vouchers can't buy lab time, detector development, or validation experiments.",
    "Vanishingly few accredited coaches have scaled particle-physics ventures.",
    "Bottom-up model may not mobilize mission-oriented physics commercialization."])

fw(("FA", "EIC · The accelerator gap"),
   "EIC: the Accelerator starts at TRL 5",
   "EIC Accelerator, Work Programme 2026.",
   ["Flagship Accelerator requires TRL 5\u20138 \u2014 years away for emerging muon ventures.",
    "Blended finance (grant up to €2.5M + equity €0.5\u201310M); success rate 3\u20137% \u2014 among the most competitive globally.",
    "Relevant early-TRL pathway: Pathfinder (TRL 1\u20134, €3\u20134M grants) + Transition (TRL 3\u20136), both needing international consortia."])

fw(("FA", "EIC · Swiss eligibility restored Jan 2025"),
   "EIC: Switzerland fully eligible since Jan 2025",
   "European Commission / SERI (Association Agreement signed Bern, 10 Nov 2025, retroactive to 1 Jan 2025).",
   ["Swiss firms fully eligible for all instruments incl. equity; can coordinate consortia.",
    "Feb 2026: 4 Swiss winners (Neutrality, Verity, Pregnolia, DePoly); June 2026: 4 of 38 \u2014 incl. BTRY (€2.2M) and Nanoflex (€12.5M).",
    "Swiss EIC winners are typically Venture Kick / Venture Leaders alumni \u2014 the national pipeline prepares them.",
    "\u2192 Sequencing: Innosuisse early \u2192 EIC Pathfinder/Transition \u2192 Accelerator \u2192 STEP Scale Up."])

fw(("FA", "Candidates · CERN Venture Connect"),
   "Candidate: CERN Venture Connect",
   "CERN Knowledge Transfer Highlights 2025.",
   ["CERN faces the IDENTICAL challenge as NCCR Muoniverse: fundamental particle physics \u2192 commercial applications.",
    "CVC model: 0% equity, express agreements, curated technology portfolio, 60 partner orgs across 18 countries.",
    "8 startups joined in 2025, collectively raising >€5M.",
    "\u2192 The most directly comparable institutional model \u2014 must be benchmarked in Part 2."])

fw(("FA", "Candidates · I-Corps"),
   "Candidate: NSF I-Corps",
   "Hayter et al. (2025), 'Launching lean', J. Technology Transfer.",
   ["The most internationally validated early-TRL commercialization program \u2014 adopted in 8 countries (EXIST Germany, Brazil, +6).",
    "Core methodology: lean startup + intensive customer discovery, structured 7-week immersion.",
    "Direct comparator to Innosuisse's bespoke coaching: methodology-driven vs mentor-driven.",
    "\u2192 STRONGLY RECOMMENDED for the November 1 benchmarking milestone."])

# ---- SECTION 3: DATA ----
S3 = "DATA COMPENDIUM"
A3 = SKY
def dt(key, title, ref, comments):
    slides.append(("dt", key, title, ref, comments))

dt(("DC", "Data · The binding constraint"),
   "Market features dominate (β = 0.514)",
   "Foresight (2024), J. Innovation & Knowledge \u2014 PLS-SEM, n = 110 EU deep tech firms, R\u00b2 = 0.47.",
   ["Market features β = 0.514 \u2014 strongest predictor by far.",
    "Technology β = 0.310, environmental β = 0.306.",
    "Firm internal features: NOT significant.",
    "\u2192 Market identification is the single highest-leverage activity."])

dt(("DC", "Data · Ecosystem: talent IRR=1.66"),
   "Talent and knowledge drive deep tech",
   "Raff-Heinen et al. (2026), J. Technology Transfer \u2014 272 regions, 1,605 startups.",
   ["Talent pools: IRR 1.66 \u2014 strongest predictor.",
    "Knowledge infrastructure: IRR 1.48.",
    "Formal institutions (policy, IP law): NOT significant.",
    "Top 10 regions = 40% of all deep tech startups."])

dt(("DC", "Data · Incubation: university r=0.118"),
   "Incubation: type and age matter",
   "Seitz et al. (2026), Review of Managerial Science \u2014 meta-analysis, 22 studies, ~37,000 ventures.",
   ["University incubators r = 0.118 vs overall 0.067; older (>20 yrs) r = 0.183; generic business incubators \u22120.011 (NS).",
    "Networking service r = 0.071; mentoring r = \u22120.012 (NS).",
    "Technology/R&D ventures benefit more (0.082\u20130.130)."])

dt(("DC", "Data · Advisors: +0.23 vs −0.31"),
   "Advisors: business +0.23, technology \u22120.31",
   "Denoo et al. (2024), J. Technology Transfer 49:1567\u20131705 \u2014 n = 112 ventures.",
   ["Each business advisor: +0.23 AR units (p = 0.04).",
    "Each technology advisor: \u22120.31 AR units (p = 0.05).",
    "4 business advisors \u2248 +1 unit \u2248 15 yrs experience; 3 tech advisors \u2248 \u22121 unit.",
    "80%+ of ventures used at least one advisor \u2014 ubiquitous policy tool."])

dt(("DC", "Data · TTO survival: p=0.0034"),
   "TTO activities that predict survival",
   "Kruachottikul et al. (2023), J. Technology Transfer \u2014 231 spin-offs, 15-yr Cox survival panel.",
   ["Business relationships p = 0.0034 \u2014 strongest.",
    "Customer introductions p = 0.006; market knowledge p = 0.015; tech knowledge p = 0.042.",
    "Personal relationships: marginal (p = 0.051).",
    "\u2192 Professionalism in relationships beats personal rapport."])

dt(("DC", "Data · Switzerland #1 (100 vs 33)"),
   "Switzerland: #1 per capita, 3\u00d7 the runner-up",
   "Deep Tech Nation Switzerland / Dealroom.co (2025), European Spinout Report.",
   ["Spinout value per capita index: 100 vs 33 (runner-up).",
    "ETH #3, EPFL #4, UZH #8 in Europe.",
    "60\u201363% of Swiss VC flows to deep tech \u2014 #1 globally (China 56%, US 54%).",
    "3 of 6 global >$1B spinout exits in 2025 were Swiss (Araris, Nexthink, u-blox)."])

dt(("DC", "Data · Innosuisse: 57% market launch"),
   "Innosuisse Innovation Projects: 57% launch rate",
   "Innosuisse Impact Monitor 2021\u20132023.",
   ["57% market launch within 3 years post-funding.",
    "2.5 highly qualified jobs per project (~800\u2013900 FTE per cohort).",
    "Leverage: 1 CHF \u2192 4 CHF value added.",
    "~50% continue with the research partner after 3 years."])

dt(("DC", "Data · BRIDGE: 70% founding rate"),
   "BRIDGE PoC: 70% founding rate",
   "econcept AG evaluation \u00b7 Innosuisse Impact Monitor 2021\u20132023.",
   ["70% of funded researchers go on to found a startup \u2014 highest conversion of any Innosuisse instrument.",
    "CHF 130K over 12 months; 80 startups created from 121 PoC projects.",
    "\u2192 Proves the concept works when the timeline fits \u2014 muon timelines do not."])

dt(("DC", "Data · KOF DiD: +20.7% sales"),
   "KOF causal effect: +20.7% sales, +17.6% employment",
   "Hulfeld, Spescha & W\u00f6rter (KOF ETH Zurich, Nov 2024) \u2014 DiD, 3,220 firm-years, 920 firms.",
   ["+20.7% sales and +17.6% employment over 5 years vs matched controls.",
    "Employment effect grows: +7.8% (yr 0) \u2192 +17.2% (yr 2) \u2192 +27.9% (yr 4).",
    "Statistically significant only for small firms (5\u201350 employees) with ETH Domain partners."])

dt(("DC", "Data · TRL coverage gap"),
   "The TRL 1\u20133 gap, visualized",
   "Milestone 1 Framework Assessment V3, \u00a710 Framework Comparison.",
   ["RCA/Imperial: TRL 4\u20136. CRL: 1\u20139 (conceptual only). Innosuisse: 4\u20137. EIC: 1\u20138 pipeline.",
    "Red zone = TRL 1\u20133 \u2014 where NCCR Muoniverse operates.",
    "Only Innosuisse has independent causal validation (KOF DiD); only EIC publishes an assessment methodology."])

dt(("DC", "Data · EU vs US: €58B vs €215B"),
   "EU vs US: a 3.7\u00d7 capital gap",
   "Stryber (2025) \u2014 interviews incl. CERN, ETH AI Center, Fraunhofer Venture.",
   ["Deep tech VC: €58B (EU) vs €215B (US) \u2014 3.7\u00d7 gap.",
    "62% of EU investors reject failed founders vs 23% US (2.7\u00d7 more risk-averse).",
    "PhD-to-startup: 15% vs 45%; corporate engagement: 12% vs 40%.",
    "US holds 49% of global deep tech VC."])

dt(("DC", "Data · Deep tech: +35% time, +48% capital"),
   "Deep tech: +35% time, +48% capital",
   "UNDP (2025) \u00b7 Kortsch et al. (2024).",
   ["Time to milestones: 135 vs 100 (+35%).",
    "Capital required: 148 vs 100 (+48%).",
    "Development timeline: 1\u20133 yrs conventional vs 5\u201315+ yrs deep tech.",
    "97% of deep tech ventures contribute to \u22651 SDG."])

# =====================================================================
# BUILD
# =====================================================================
# Title slide
s = add_slide()
wallet_bar(s, "MILESTONE 1", TEAL, 0, total)
textbox(s, 0.55, 2.4, 12.2, 0.5, [("NCCR MUONIVERSE \u00b7 BACHELOR THESIS \u00b7 HSG", 13, True, TEAL, 0)])
textbox(s, 0.55, 2.8, 12.2, 1.6, [
    ("Milestone 1 Findings", 44, True, WHITE, 4),
    ("Literature Review \u00b7 Framework Assessment \u00b7 Data Compendium", 18, False, SLATE, 0)])
textbox(s, 0.55, 4.6, 12.2, 0.8, [
    ("One page per finding: document screenshot with highlighted quote, source reference, and commentary.", 13, False, GRAY, 3),
    ("Lorik Dalloshi \u00b7 prepared for the NCCR Muoniverse team \u00b7 August 2026", 11, False, GRAY, 0)])

# Section dividers + findings
p = 2
divider("LITERATURE REVIEW", TEAL, "Seven themes", "What the 53 sources converge on", p, total); p += 1
for kind, key, title, ref, comments in slides:
    if kind == "theme":
        finding_slide("LITERATURE REVIEW", TEAL, title, "Literature Review \u00b7 Finding", key, ref, comments, p, total); p += 1

divider("FRAMEWORK ASSESSMENT", VIOLET, "Four frameworks, one gap", "RCA/Imperial \u00b7 CRL \u00b7 Innosuisse \u00b7 EIC \u00b7 + 4 candidates", p, total); p += 1
for kind, key, title, ref, comments in slides:
    if kind == "fw":
        finding_slide("FRAMEWORK ASSESSMENT", VIOLET, title, "Framework Assessment \u00b7 Finding", key, ref, comments, p, total); p += 1

divider("DATA COMPENDIUM", SKY, "The numbers", "Statistical evidence behind every claim", p, total); p += 1
for kind, key, title, ref, comments in slides:
    if kind == "dt":
        finding_slide("DATA COMPENDIUM", SKY, title, "Data Compendium \u00b7 Finding", key, ref, comments, p, total); p += 1

# Closing slide
s = add_slide()
wallet_bar(s, "MILESTONE 1", TEAL, p, total)
textbox(s, 0.55, 2.9, 12.2, 1.0, [("37 findings \u00b7 1 gap", 34, True, WHITE, 0)])
textbox(s, 0.55, 3.9, 12.2, 0.9, [
    ("The headline: no existing framework was designed for TRL 1\u20133 \u2014 where NCCR Muoniverse operates.", 14, False, SLATE, 4),
    ("Next: interview phase (Milestone 2, Sept 1) \u00b7 benchmarking (Milestone 3, Nov 1).", 12, False, GRAY, 0)])

out = r"G:\Mon Drive\ZHermes Agent\Bachelor Thesis\presentations\Milestone1_Findings_Deck.pptx"
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("SAVED:", out, os.path.getsize(out), "bytes,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
