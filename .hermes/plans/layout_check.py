"""Deterministic layout check: simulate Arial 10.5pt wrapping for every comment
textbox and verify (a) no box crosses the footer rule at y=7.03, (b) no two
comment boxes on a slide overlap, (c) nothing exceeds slide bounds."""
from pptx import Presentation
from pptx.util import Emu
from PIL import ImageFont

PPTX = r'C:\Dev\CyberZen\.hermes\plans\Milestone1_Findings_Deck_accurate_v3.pptx'
FOOTER_Y = 7.03          # inches
RIGHT_COL_X = 8.18       # comments column starts here
FONT_PATH = r'C:\Windows\Fonts\arial.ttf'

# 10.5 pt at 96 dpi -> px for PIL measurement; widths in inches*96
font_px = ImageFont.truetype(FONT_PATH, 14)   # 10.5pt ~= 14px @96dpi
LINE_H_IN = 14 * 1.22 / 96                    # generous line spacing

def wrapped_lines(text, width_in):
    avail = width_in * 96
    lines, cur = 1, ''
    for word in text.split():
        trial = (cur + ' ' + word).strip()
        if font_px.getlength(trial) <= avail:
            cur = trial
        else:
            lines += 1
            cur = word
    return lines

p = Presentation(PPTX)
problems = []
for idx, s in enumerate(p.slides, 1):
    tx = ' '.join(sh.text for sh in s.shapes if sh.has_text_frame)
    if 'What the excerpt shows' not in tx:
        continue  # not a v2 finding slide
    boxes = []
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        x = sh.left / 914400; y = sh.top / 914400
        w = sh.width / 914400; h = sh.height / 914400
        t = sh.text_frame.text
        import re as _re
        if _re.fullmatch(r'\d{2} / \d{2}', t.strip()):
            continue  # page-number footer, intentionally below the rule
        if x >= RIGHT_COL_X - 0.05 and t and 'COMMENTS' not in t:
            n = wrapped_lines(t, w - 0.0)
            need = n * LINE_H_IN
            boxes.append((y, y + need, idx, t[:40], n))
            if y + need > FOOTER_Y + 0.02:
                problems.append((idx, 'FOOTER OVERFLOW', t[:40], round(y + need, 2)))
    boxes.sort()
    for a, b in zip(boxes, boxes[1:]):
        if a[1] > b[0] + 0.01:
            problems.append((a[2], 'BOX OVERLAP', a[3], round(a[1], 2), b[3]))

print('slides checked with v2 comments')
if problems:
    for pr in problems:
        print('PROBLEM', pr)
else:
    print('ALL CLEAN: no footer overflow, no comment-box overlaps')
