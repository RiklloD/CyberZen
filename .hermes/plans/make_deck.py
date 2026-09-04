#!/usr/bin/env python3
"""Generate the Milestone 2 Interview Plan deck (PPTX) for Shaun/Bill."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------- palette ----------
NAVY   = RGBColor(0x0F, 0x1B, 0x2D)   # deep navy
TEAL   = RGBColor(0x14, 0xB8, 0xA6)   # teal accent
LIGHT  = RGBColor(0xF5, 0xF8, 0xFA)   # light bg
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SLATE  = RGBColor(0x33, 0x41, 0x55)   # body text
GRAY   = RGBColor(0x8A, 0x94, 0xA6)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)

FONT = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s

def textbox(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color, space_after)"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for text, size, bold, color, space in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run(); r.text = text
        f = r.font
        f.name = FONT; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    return tb

def chip(s, x, y, w, h, text, fill=TEAL, txt=WHITE, size=12, bold=True):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = txt
    return c

def header(s, kicker, title, page):
    textbox(s, 0.6, 0.35, 9.5, 0.4, [(kicker, 12, True, TEAL, 0)])
    textbox(s, 0.6, 0.7, 11.5, 0.8, [(title, 28, True, NAVY, 0)])
    # accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(2.2), Inches(0.045))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background(); bar.shadow.inherit = False
    textbox(s, 12.2, 0.35, 0.8, 0.4, [(f"{page:02d}", 12, True, GRAY, 0)], align=PP_ALIGN.RIGHT)

def footer(s, note):
    textbox(s, 0.6, 7.05, 12.1, 0.35, [(note, 9, False, GRAY, 0)])

# ============================================================ 1 TITLE
s = add_slide(NAVY)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height)
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background(); bar.shadow.inherit = False
textbox(s, 1.0, 1.7, 11.3, 0.5, [("NCCR MUONIVERSE  ·  BACHELOR THESIS  ·  HSG", 14, True, TEAL, 0)])
textbox(s, 1.0, 2.2, 11.3, 1.6, [
    ("Interview Phase Plan", 44, True, WHITE, 6),
    ("Deep Tech Commercialization at Early TRL in the Swiss Innovation Ecosystem", 20, False, RGBColor(0xC9,0xD6,0xE3), 0)])
textbox(s, 1.0, 4.3, 11.3, 1.2, [
    ("Milestone 2 · September 1, 2026", 16, True, WHITE, 4),
    ("Lorik Dalloshi — with Shaun & Bill (NCCR Muoniverse)", 13, False, RGBColor(0x9F,0xB3,0xC8), 0)])
textbox(s, 1.0, 6.3, 11.3, 0.5, [("Draft v1.0 · August 9, 2026 · for review", 11, False, RGBColor(0x6B,0x7F,0x96), 0)])

# ============================================================ 2 WHY INTERVIEWS
s = add_slide()
header(s, "CONTEXT", "Why interviews — and why now", 2)
textbox(s, 0.6, 1.85, 6.0, 4.6, [
    ("The thesis has four parts. Milestone 1 (July) delivered the literature map and framework assessment. Milestone 2 (September 1) is the interview phase.", 14, False, SLATE, 10),
    ("The literature tells us what support mechanisms exist and what the evidence says. It cannot tell us what actually happens inside a physics venture at TRL 1–3 — that only exists in the heads of the people who lived it.", 14, False, SLATE, 10),
    ("Interviews fill the gap between theory and practice: they test whether frameworks guide decisions or are just paperwork, and they collect the raw material for the thesis's proposed framework adaptation.", 14, False, SLATE, 0)])
# right panel: the 4 parts
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.85), Inches(5.7), Inches(4.4))
panel.fill.solid(); panel.fill.fore_color.rgb = NAVY; panel.line.fill.background(); panel.shadow.inherit = False
textbox(s, 7.35, 2.05, 5.0, 0.4, [("Thesis structure", 15, True, WHITE, 0)])
parts = [
    ("PART 1", "Literature review", "53 sources · 7 themes · 8 gaps", "DONE — July 2026"),
    ("PART 2", "Benchmarking", "4 frameworks scored · 4 shortlisted", "NEXT — Nov 1"),
    ("PART 3", "Firm interviews", "6–10 Swiss experts · 7 questions", "NOW — Sep 1"),
    ("PART 4", "Synthesis", "Convergence matrix · final thesis", "May 2027"),
]
y = 2.6
for tag, t, sub, status in parts:
    chip(s, 7.35, y, 1.15, 0.5, tag, fill=TEAL, size=10)
    textbox(s, 8.65, y-0.06, 3.9, 0.6, [(t, 13, True, WHITE, 0)])
    textbox(s, 8.65, y+0.28, 3.9, 0.4, [(sub, 10, False, RGBColor(0x9F,0xB3,0xC8), 0)])
    textbox(s, 7.35, y+0.62, 5.0, 0.35, [(status, 10, True, AMBER, 0)])
    y += 1.02
footer(s, "Source: PLAN.md & MASTER_TODO.md (Bachelor Thesis workspace)")

# ============================================================ 3 THE 7 QUESTIONS
s = add_slide()
header(s, "CORE DELIVERABLE", "The 7 interview questions", 3)
textbox(s, 0.6, 1.7, 12.1, 0.4, [
    ("Each question is grounded in the literature (Milestone 1) and mapped to a research gap. Full guide with probes, rationale and codes: interviews/interview_guide.md", 12, False, GRAY, 0)])
qs = [
    ("1", "Your journey", "Lab to today — what drove the decision to commercialize?", "Motivation · TRL trajectory"),
    ("2", "Support mechanisms", "Which actually helped? Which were just paperwork?", "Theme 4 · Denoo 2024"),
    ("3", "Frameworks", "Did TRL/CRL/I-Corps guide decisions — or reporting only?", "Benchmarking core"),
    ("4", "The early-TRL gap", "What was missing at TRL 1–3? What would have helped?", "Central thesis gap"),
    ("5", "Swiss ecosystem", "What's unique about commercializing physics in CH?", "Gap 1 · TTO study"),
    ("6", "Market discovery", "How do you find a market when none is obvious?", "Universe of markets"),
    ("7", "Open reflection", "If you built the support program from scratch, what would it look like?", "Design input"),
]
x0, y0, cw, ch, gx, gy = 0.6, 2.3, 3.9, 1.28, 0.15, 0.15
for i, (n, t, q, why) in enumerate(qs):
    x = x0 + (i % 3) * (cw + gx)
    y = y0 + (i // 3) * (ch + gy)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(ch))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0); card.line.width = Pt(1)
    card.shadow.inherit = False
    textbox(s, x+0.2, y+0.12, 0.5, 0.4, [(n, 18, True, TEAL, 0)])
    textbox(s, x+0.65, y+0.16, cw-0.85, 0.4, [(t, 13, True, NAVY, 0)])
    textbox(s, x+0.2, y+0.55, cw-0.4, 0.5, [(q, 10.5, False, SLATE, 0)])
    textbox(s, x+0.2, y+1.0, cw-0.4, 0.3, [(why, 8.5, True, GRAY, 0)])
footer(s, "60-minute structure per Shaun's call of 08.08.2026: Background 5' · Journey 15' · Frameworks 15' · Gap 10' · Ecosystem 10' · Reflection 5'")

# ============================================================ 4 WHO WE INTERVIEW
s = add_slide()
header(s, "CORE DELIVERABLE", "Who we interview — 16 named candidates", 4)
textbox(s, 0.6, 1.7, 12.1, 0.4, [
    ("Tiered list, verified via public sources (Aug 2026). Target: 6–10 completed interviews. Full list with contact paths: interviews/contact_list.md", 12, False, GRAY, 0)])
tiers = [
    ("TIER 1", "TTOs & funders", "John Millard (PSI TTO) · Andrea Crottini (EPFL TTO) · Effy Vayena (ETH) · Han Dols (CERN KT) · Christoph Falk (Innosuisse) · Christian Brunner (BRIDGE)", TEAL),
    ("TIER 2", "Muon & physics owners", "Marc Janoschek (PSI muon center) · Klaus Kirch (NCCR director) · Matthias Schneebeli (DECTRIS CEO) · Christian Brönnimann (DECTRIS founder)", NAVY),
    ("TIER 3", "Ecosystem builders", "Beat Schillig (Venture Kick) · Philip Hassler (Venture Kick) · Dominique Gruhl-Bégin (Innosuisse CEO)", TEAL),
    ("TIER 4", "Academics (Shaun's picks)", "Henry Chesbrough (open innovation) · Christopher Tucci (Imperial, ex-EPFL) · Dale Rickert (Handelsblatt TECH)", NAVY),
    ("TIER 5", "Venture founders", "3 slots — names via Shaun/Bill (PSI & ETH/EPFL physics spin-offs, Innosuisse-funded deep tech)", TEAL),
]
y = 2.25
for tag, t, names, col in tiers:
    chip(s, 0.6, y, 1.0, 0.55, tag, fill=col, size=10)
    textbox(s, 1.8, y-0.02, 2.6, 0.6, [(t, 12.5, True, NAVY, 0)])
    textbox(s, 4.5, y-0.02, 8.2, 0.6, [(names, 11, False, SLATE, 0)])
    y += 0.78
textbox(s, 0.6, 6.35, 12.1, 0.5, [
    ("Priority order: Tier 1 + Tier 2 first (most thesis-relevant, most likely to accept via NCCR intros). Shaun/Bill intros are the single biggest acceptance lever.", 11.5, True, AMBER, 0)])
footer(s, "All Tier 1–3 contacts have public contact paths; Tier 4 via LinkedIn (from Shaun's call notes)")

# ============================================================ 5 OUTREACH & TIMELINE
s = add_slide()
header(s, "PLAN", "Outreach & timeline — August to December", 5)
# timeline bar
y = 2.1
textbox(s, 0.6, y, 12.1, 0.4, [("Interview phase (overlaps benchmarking by design — PLAN.md)", 12, False, GRAY, 0)])
steps = [
    ("AUG 10–22", "Review with Shaun/Bill", "Validate guide, contact list, ethics check", "Shaun/Bill"),
    ("AUG 24–SEP 5", "First outreach wave", "Tier 1+2 emails (templates ready), intros requested", "Lorik"),
    ("SEP 1", "MILESTONE 2", "Interview plan + guide delivered", "Lorik → Shaun"),
    ("SEP 8–OCT 30", "Conduct interviews", "6–10 × 45–60 min, memo within 24h, transcribe", "Lorik"),
    ("NOV 1", "MILESTONE 3", "Benchmarking (overlaps)", "Lorik → Shaun"),
    ("NOV–DEC", "Coding & analysis", "Thematic coding, evidence table, Part 3 draft", "Lorik"),
]
y = 2.6
for when, what, how, who in steps:
    chip(s, 0.6, y, 1.7, 0.5, when, fill=NAVY, size=9.5)
    textbox(s, 2.5, y-0.03, 3.2, 0.6, [(what, 12.5, True, NAVY, 0)])
    textbox(s, 5.8, y-0.03, 5.2, 0.6, [(how, 10.5, False, SLATE, 0)])
    textbox(s, 11.1, y-0.03, 1.7, 0.6, [(who, 10, True, TEAL, 0)])
    y += 0.72
textbox(s, 0.6, 6.5, 12.1, 0.5, [
    ("Risk: acceptance rate ~50% → 16 candidates for 6–10 interviews. If Tier 1 stalls, Tier 3–5 fill the quota. Ethics approval (if required) must be secured before first outreach.", 11, False, GRAY, 0)])
footer(s, "Timeline per PLAN.md Phase 3 (Weeks 10–18)")

# ============================================================ 6 ETHICS & CONSENT
s = add_slide()
header(s, "COMPLIANCE", "Ethics & data handling", 6)
cards = [
    ("Consent first", "Information sheet + consent form sent 48h before each interview. Three anonymization options: fully anonymized / org named / fully named.", TEAL),
    ("Recording", "Audio recording only with explicit consent (written + verbal re-confirmation on tape). Otherwise detailed notes.", NAVY),
    ("Storage", "Transcripts stored locally, password-protected. Deleted after thesis submission (May 2027).", TEAL),
    ("Withdrawal", "Participants may withdraw at any time, including after the interview. Transcripts shared for review and correction.", NAVY),
]
x0, y0, cw, ch, gx = 0.6, 2.0, 6.0, 1.9, 0.25
for i, (t, d, col) in enumerate(cards):
    x = x0 + (i % 2) * (cw + gx)
    y = y0 + (i // 2) * (ch + 0.25)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(ch))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0); card.line.width = Pt(1)
    card.shadow.inherit = False
    chip(s, x+0.25, y+0.25, 1.6, 0.4, t, fill=col, size=10)
    textbox(s, x+0.25, y+0.85, cw-0.5, 0.95, [(d, 11, False, SLATE, 0)])
textbox(s, 0.6, 6.3, 12.1, 0.6, [
    ("Open question for Shaun: does HSG require formal ethics approval for semi-structured interviews with professionals? If yes, templates are ready for submission.", 12, True, AMBER, 0)])
footer(s, "Templates: interviews/consent_template.md")

# ============================================================ 7 DELIVERABLES
s = add_slide()
header(s, "OUTPUTS", "What this phase produces", 7)
dels = [
    ("Interview guide", "7 questions + probes + literature rationale + codes", "interviews/interview_guide.md", "DONE"),
    ("Contact list", "16 named candidates in 5 tiers, with contact paths", "interviews/contact_list.md", "DONE"),
    ("Outreach kit", "4 email templates (cold, intro, follow-up, thank-you)", "interviews/outreach_email_templates.md", "DONE"),
    ("Consent kit", "Information sheet + consent form + verbal script", "interviews/consent_template.md", "DONE"),
    ("Memo template", "Post-interview memo within 24h (surprises, quotes, codes)", "interviews/memo_template.md", "DONE"),
    ("Coding framework", "Deductive codes from literature + inductive log", "interviews/coding/coding_framework.md", "DONE"),
    ("Part 3 draft", "8–10 pages thematic analysis (after interviews)", "drafts/V3_interviews.md", "NOV–DEC"),
    ("Evidence table", "Anonymized interviewee × theme × quote matrix", "interviews/coding/evidence_table.md", "NOV–DEC"),
]
y = 1.95
for t, d, path, status in dels:
    chip(s, 0.6, y, 1.5, 0.45, status, fill=TEAL if status == "DONE" else AMBER, size=9)
    textbox(s, 2.3, y-0.04, 3.4, 0.5, [(t, 12.5, True, NAVY, 0)])
    textbox(s, 5.8, y-0.04, 5.3, 0.5, [(d, 10.5, False, SLATE, 0)])
    textbox(s, 11.2, y-0.04, 1.7, 0.5, [(path, 8.5, False, GRAY, 0)])
    y += 0.62
footer(s, "All files live in the Bachelor Thesis workspace (G:\\Mon Drive\\ZHermes Agent\\Bachelor Thesis)")

# ============================================================ 8 OPEN QUESTIONS
s = add_slide()
header(s, "FOR SHAUN & BILL", "Decisions needed before first outreach", 8)
qs = [
    ("1", "Does the 7-question structure match what you expect? Strict guide or loose topic list?"),
    ("2", "Can you provide intros to John Millard (PSI TTO), Marc Janoschek (PSI muon), DECTRIS, EPFL TTO?"),
    ("3", "Is HSG ethics approval required? (Templates ready either way.)"),
    ("4", "Tier 5 founders: which PSI/ETH/EPFL physics spin-offs should we target? Names please."),
    ("5", "What did you mean by 'Chipeno' and 'Api 7' in the call notes?"),
    ("6", "Is the Springer chapter (10.1007/978-3-031-86958-7_10, Smart Services Summit) the intended reference?"),
]
y = 2.0
for n, q in qs:
    chip(s, 0.6, y, 0.5, 0.5, n, fill=NAVY, size=12)
    textbox(s, 1.35, y-0.02, 11.3, 0.6, [(q, 13, False, SLATE, 0)])
    y += 0.72
textbox(s, 0.6, 6.4, 12.1, 0.5, [
    ("Next step: your feedback → finalize guide → ethics check → first outreach wave (Aug 24).", 12.5, True, TEAL, 0)])
footer(s, "Full context: notes/open_questions.md · Feedback from calls/Call 08.08.26.docx")

# ============================================================ 9 THANK YOU
s = add_slide(NAVY)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height)
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background(); bar.shadow.inherit = False
textbox(s, 1.0, 2.6, 11.3, 1.0, [("Thank you", 40, True, WHITE, 0)])
textbox(s, 1.0, 3.6, 11.3, 1.2, [
    ("Feedback welcome — the guide is a draft until you say otherwise.", 16, False, RGBColor(0xC9,0xD6,0xE3), 6),
    ("Lorik Dalloshi · HSG · NCCR Muoniverse", 13, False, RGBColor(0x9F,0xB3,0xC8), 0)])

out = r"G:\Mon Drive\ZHermes Agent\Bachelor Thesis\interviews\Interview_Plan_Milestone2.pptx"
prs.save(out)
print("SAVED:", out, os.path.getsize(out), "bytes")
